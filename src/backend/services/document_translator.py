"""Document translation loop (Sprint 3 + 7, SPEC §3.4 / §4.2).

Three functions behind one per-format interface:

    extract(path) -> IR            list of segments (+ format), granularity-preserving
    translate(IR, ...) -> IR       per-language, per-segment two-pass fill
    recompose(IR, lang, out) -> path   rebuild the document in the original format

``translate`` is format-agnostic — it only touches each segment's ``text`` /
``translate`` / ``out``. Only ``extract`` and ``recompose`` are per-format:

  - **Tier 1** (`.txt` / `.md`): loss-less line segmentation (blank runs + fenced
    code are skeleton), so recomposition is byte-stable.
  - **Tier 2** (`.docx` / `.rtf`, Sprint 7): in-container libs (lesson #11) — no
    host CLI. `.docx` via **python-docx** at paragraph granularity (paragraph
    style kept; minor inline-format shifts accepted, §4.2); `.rtf` via **pandoc**
    (rtf↔markdown), reusing the Tier-1 markdown segmenter.
  - **Tier 3** (`.pptx`, Sprint 8, experimental/post-MVP — lesson #12): via
    **python-pptx**, text runs per shape at paragraph granularity. Best-effort;
    complex layouts may degrade. The UI warns before submission (§10).

Each segment is translated two-pass (draft with neighbour context → glossary
review) through the connection registry — never a hardcoded endpoint (lesson #9).
"""

import logging
import re
from pathlib import Path
from typing import Any

from src.api.v1.admin.llm import (
    get_llm_settings, load_document_translation_prompt, load_plain_translation_prompt,
    load_document_flavour, load_review_prompt, load_context_prompt,
)
from src.core.config import config
from src.core.languages import language_name
from src.services.glossary_slice import slice_glossary, format_glossary_block
from src.services.llm_provider import llm

logger = logging.getLogger("backend.document_translator")

_MD_SUFFIXES = {".md", ".markdown"}
_DOCX_SUFFIXES = {".docx"}
_RTF_SUFFIXES = {".rtf"}
_PPTX_SUFFIXES = {".pptx"}

# RTL text direction (ADR-020). Of the 17 languages only Arabic and Urdu are
# right-to-left; recompose sets RTL direction for these targets in docx/pptx.
_RTL_LANGS = {"ar", "ur"}


def _is_rtl(lang: str) -> bool:
    return (lang or "").split("-")[0].lower() in _RTL_LANGS


def _has_rtl_char(text: str) -> bool:
    """True if any character is Arabic-script (covers AR + UR). Used to mark only
    RTL-script runs `w:rtl`, leaving Latin/number runs LTR (as Word itself does)."""
    # Arabic + Supplement + Extended-A + Presentation Forms A/B.
    return any(
        0x0600 <= o <= 0x06FF or 0x0750 <= o <= 0x077F or 0x08A0 <= o <= 0x08FF
        or 0xFB50 <= o <= 0xFDFF or 0xFE70 <= o <= 0xFEFF
        for o in map(ord, text)
    )


def _strip_think(text: str) -> str:
    """Remove <think>...</think> reasoning blocks from model output."""
    return re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()


# Completeness guard (queso-suizo bug): the model occasionally returns an empty or
# early-EOS'd unit — empty after <think> stripping, a premature stop, or a chunk lost
# to model eviction. Nothing used to check, so the hole was silently concatenated and
# emailed. We now retry, and fail the job loudly rather than ship a truncated document.
_TRANSLATION_MAX_ATTEMPTS = 3
_MIN_OUTPUT_RATIO = 0.4       # a prose unit shorter than this fraction of its source is truncated/empty
_MIN_GUARDED_CHARS = 200      # don't length-check tiny units (short paragraphs, headings, cells)


def _looks_truncated(core: str, out: str) -> bool:
    """True if a translation is empty or suspiciously shorter than its source — the
    signature of an emptied/early-stopped unit that must NOT be shipped silently.
    Only length-checks non-trivial units, so short structural paragraphs (which can
    legitimately compress) never false-positive."""
    stripped = out.strip()
    if not stripped:
        return True
    src = core.strip()
    return len(src) >= _MIN_GUARDED_CHARS and len(stripped) < _MIN_OUTPUT_RATIO * len(src)


# ---------------------------------------------------------------------------
# Segmentation (Tier 1: txt / md) — loss-less
# ---------------------------------------------------------------------------


def _segment(text: str, fmt: str) -> list[dict[str, Any]]:
    """Split into ordered segments, each ``{text, translate}``.

    Three kinds, so structure survives recomposition:
      - blank runs (one or more blank lines) → skeleton, not translated;
      - fenced code blocks (```) in markdown → verbatim, not translated;
      - runs of consecutive non-blank lines → one translatable text segment.

    Loss-less: ``"".join(s["text"] for s in _segment(t, fmt)) == t``.
    """
    lines = text.splitlines(keepends=True)
    segments: list[dict[str, Any]] = []
    i, n = 0, len(lines)
    while i < n:
        stripped = lines[i].strip()
        # fenced code block (md only) — consume until the closing fence
        if fmt == "md" and stripped.startswith("```"):
            block = [lines[i]]
            j = i + 1
            while j < n:
                block.append(lines[j])
                closed = lines[j].strip().startswith("```")
                j += 1
                if closed:
                    break
            segments.append({"text": "".join(block), "translate": False})
            i = j
            continue
        # blank run
        if stripped == "":
            j = i
            while j < n and lines[j].strip() == "":
                j += 1
            segments.append({"text": "".join(lines[i:j]), "translate": False})
            i = j
            continue
        # text run — consecutive non-blank, non-fence lines
        j = i
        while j < n:
            s = lines[j].strip()
            if s == "" or (fmt == "md" and s.startswith("```")):
                break
            j += 1
        segments.append({"text": "".join(lines[i:j]), "translate": True})
        i = j
    return segments


# ---------------------------------------------------------------------------
# Tier 2 handlers (docx / rtf) — in-container libs only (lesson #11)
# ---------------------------------------------------------------------------


# docx fidelity (run-level formatting preservation). We edit a *copy* of the
# source, so structure (tables, images, lists, numbering, headings, borders,
# hyperlinks) is preserved for free. The only thing lost by the old "dump all
# text in run[0]" was **inline run formatting** (bold/italic/highlight/color/size)
# when a paragraph mixed formats. So we translate each paragraph carrying inline
# **markers** ⟦i⟧…⟦/i⟧ around each formatting span, and on recompose we put each
# translated span back into its own run (keeping that run's full rPr). Runs that
# carry a drawing/field are opaque (kept verbatim → inline images survive).
_MARK_RE = re.compile(r"⟦(\d+)⟧(.*?)⟦/\1⟧", re.DOTALL)
_XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


def _run_text(r: Any, qn: Any) -> str:
    return "".join(t.text or "" for t in r.findall(qn("w:t")))


def _run_is_opaque(r: Any, qn: Any) -> bool:
    """A run carrying an image/field/object — keep it verbatim, never translate."""
    return any(r.find(qn(t)) is not None for t in ("w:drawing", "w:pict", "w:object", "w:fldChar"))


def _rpr_sig(r: Any, qn: Any) -> str:
    rpr = r.find(qn("w:rPr"))
    return rpr.xml if rpr is not None else ""


def _p_spans(p: Any, qn: Any) -> list[dict[str, Any]]:
    """Ordered spans of a paragraph (direct children). Consecutive `w:r` runs with
    the same rPr merge into one text span; a run with a drawing/field is its own
    opaque span; each `w:hyperlink` is one translatable span (its link is kept)."""
    spans: list[dict[str, Any]] = []
    for child in p:
        tag = child.tag
        if tag == qn("w:r"):
            if _run_is_opaque(child, qn):
                spans.append({"kind": "opaque", "runs": [child]})
            else:
                sig = _rpr_sig(child, qn)
                if spans and spans[-1]["kind"] == "text" and spans[-1]["sig"] == sig:
                    spans[-1]["runs"].append(child)
                else:
                    spans.append({"kind": "text", "sig": sig, "runs": [child]})
        elif tag == qn("w:hyperlink"):
            spans.append({"kind": "link", "runs": child.findall(qn("w:r"))})
    return spans


def _span_text(span: dict[str, Any], qn: Any) -> str:
    return "".join(_run_text(r, qn) for r in span["runs"])


def _set_run_text(r: Any, text: str, qn: Any, OxmlElement: Any) -> None:
    """Replace a run's textual content with one <w:t>, keeping its rPr."""
    for tag in ("w:t", "w:br", "w:tab", "w:cr"):
        for el in r.findall(qn(tag)):
            r.remove(el)
    t = OxmlElement("w:t")
    t.set(_XML_SPACE, "preserve")
    t.text = text
    r.append(t)


def _parse_markers(text: str, n: int) -> list[str] | None:
    """Recover the n translated spans from ⟦i⟧…⟦/i⟧; None if any is missing."""
    found: dict[int, str] = {}
    for m in _MARK_RE.finditer(text):
        i = int(m.group(1))
        if 0 <= i < n and i not in found:
            found[i] = m.group(2)
    return [found[i] for i in range(n)] if len(found) == n else None


def _strip_markers(text: str) -> str:
    return re.sub(r"⟦/?\d+⟧", "", _MARK_RE.sub(lambda m: m.group(2), text))


def _serialize_p(p: Any, qn: Any) -> tuple[str, bool]:
    """Return (text-for-translation, translatable). One translatable span → plain
    text; several → each wrapped in ⟦i⟧…⟦/i⟧ so formatting boundaries survive."""
    spans = _p_spans(p, qn)
    trans = [s for s in spans if s["kind"] in ("text", "link")]
    texts = [_span_text(s, qn) for s in trans]
    full = "".join(texts)
    if not full.strip():
        return full, False
    if len(trans) <= 1:
        return full, True
    return "".join(f"⟦{i}⟧{t}⟦/{i}⟧" for i, t in enumerate(texts)), True


def _rebuild_p(p: Any, translated: str, qn: Any, OxmlElement: Any) -> None:
    """Put the translated text back, one run per formatting span (each keeps its
    own rPr). If the markers came back broken, fall back to all-text-in-first-run
    (never worse than the old behaviour)."""
    trans = [s for s in _p_spans(p, qn) if s["kind"] in ("text", "link")]
    if not trans:
        return
    if len(trans) == 1:
        texts: list[str] = [translated]
    else:
        parsed = _parse_markers(translated, len(trans))
        texts = parsed if parsed is not None else [_strip_markers(translated)] + [""] * (len(trans) - 1)
    for span, text in zip(trans, texts):
        runs = span["runs"]
        if not runs:
            continue
        _set_run_text(runs[0], text, qn, OxmlElement)  # first run keeps rPr (and, for a link, its hyperlink)
        for extra in runs[1:]:
            extra.getparent().remove(extra)


def _note_paragraph_elements(doc: Any, qn: Any) -> list[Any]:
    """`<w:p>` elements inside the footnotes/endnotes parts (skipping the
    separator/continuation notes). Only parts whose element tree is mutable (and
    thus saved by python-docx) are used; others are skipped gracefully."""
    ps: list[Any] = []
    for rel in list(doc.part.rels.values()):
        rt = getattr(rel, "reltype", "")
        if getattr(rel, "is_external", False):
            continue
        if not (rt.endswith("/footnotes") or rt.endswith("/endnotes")):
            continue
        root = getattr(rel.target_part, "element", None) or getattr(rel.target_part, "_element", None)
        if root is None:
            continue
        for note in list(root.findall(qn("w:footnote"))) + list(root.findall(qn("w:endnote"))):
            if note.get(qn("w:type")) in ("separator", "continuationSeparator"):
                continue
            ps.extend(note.iter(qn("w:p")))
    return ps


def _iter_doc_paragraphs(doc: Any, include_notes: bool = False) -> list[Any]:
    """Every translatable `<w:p>` element, in a stable order: body + tables (any
    depth, via `body.iter`) then every section's headers/footers (all types), and
    — when ``include_notes`` — footnotes/endnotes. extract + recompose walk this
    identically (same doc, same flag) so segments line up."""
    from docx.oxml.ns import qn
    ps: list[Any] = list(doc.element.body.iter(qn("w:p")))
    hf_attrs = ("header", "footer", "first_page_header", "first_page_footer",
                "even_page_header", "even_page_footer")
    for section in doc.sections:
        for attr in hf_attrs:
            hf = getattr(section, attr, None)
            if hf is None:
                continue
            try:
                if getattr(hf, "is_linked_to_previous", False):
                    continue
                for para in hf.paragraphs:
                    ps.append(para._p)
                for tbl in hf.tables:
                    ps.extend(tbl._tbl.iter(qn("w:p")))
            except Exception:
                continue
    if include_notes:
        try:
            ps.extend(_note_paragraph_elements(doc, qn))
        except Exception as e:
            logger.warning(f"footnotes/endnotes extraction skipped: {e}")
    return ps


def _extract_docx(p: Path, include_notes: bool = False) -> dict[str, Any]:
    from docx import Document
    from docx.oxml.ns import qn
    doc = Document(str(p))
    segments = []
    for pel in _iter_doc_paragraphs(doc, include_notes):
        text, translate = _serialize_p(pel, qn)
        segments.append({"text": text, "translate": translate})
    return {"format": "docx", "source_path": str(p), "segments": segments, "include_notes": include_notes}


# RTL (ADR-020). w:bidi must precede these successors in pPr / sectPr child order
# (ISO 29500); insert_element_before keeps the tree schema-valid.
_PPR_BIDI_SUCCESSORS = ("w:adjustRightInd", "w:snapToGrid", "w:spacing", "w:ind",
                        "w:contextualSpacing", "w:mirrorIndents", "w:suppressOverlap",
                        "w:jc", "w:textDirection", "w:textAlignment", "w:textboxTightWrap",
                        "w:outlineLvl", "w:divId", "w:cnfStyle", "w:rPr", "w:sectPr",
                        "w:pPrChange")
_SECTPR_BIDI_SUCCESSORS = ("w:rtlGutter", "w:docGrid", "w:printerSettings", "w:sectPrChange")


def _docx_paragraph_rtl(pel: Any, qn: Any, OxmlElement: Any) -> None:
    """Right-to-left a translated `<w:p>`: `w:bidi` on its pPr + `w:rtl` on every
    run whose text is RTL script (Latin/number runs stay LTR, as Word does)."""
    pPr = pel.get_or_add_pPr()
    if pPr.find(qn("w:bidi")) is None:
        pPr.insert_element_before(OxmlElement("w:bidi"), *_PPR_BIDI_SUCCESSORS)
    for r in pel.iter(qn("w:r")):
        if not _has_rtl_char("".join(t.text or "" for t in r.findall(qn("w:t")))):
            continue
        rPr = r.get_or_add_rPr()
        if rPr.find(qn("w:rtl")) is None:
            rPr.get_or_add_rtl()


def _docx_sections_rtl(doc: Any, qn: Any, OxmlElement: Any) -> None:
    """`w:bidi` on every section's sectPr so the page base direction is RTL."""
    for section in doc.sections:
        sectPr = section._sectPr
        if sectPr.find(qn("w:bidi")) is None:
            sectPr.insert_element_before(OxmlElement("w:bidi"), *_SECTPR_BIDI_SUCCESSORS)


def _recompose_docx(ir: dict[str, Any], target_lang: str, out_path: str | None) -> str:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    doc = Document(ir["source_path"])  # fresh copy per language (no cross-lang mutation)
    rtl = _is_rtl(target_lang)
    for seg, pel in zip(ir["segments"], _iter_doc_paragraphs(doc, ir.get("include_notes", False))):
        if not seg["translate"]:
            continue
        translated = seg.get("out", {}).get(target_lang)
        if translated is None:
            continue
        try:
            _rebuild_p(pel, translated, qn, OxmlElement)
        except Exception as e:
            logger.warning(f"docx paragraph rebuild failed ({e}); plain fallback")
            try:
                _rebuild_p(pel, _strip_markers(translated), qn, OxmlElement)
            except Exception:
                pass
        if rtl:
            try:
                _docx_paragraph_rtl(pel, qn, OxmlElement)
            except Exception as e:
                logger.warning(f"docx paragraph RTL failed ({e})")
    if rtl:
        try:
            _docx_sections_rtl(doc, qn, OxmlElement)
        except Exception as e:
            logger.warning(f"docx section RTL failed ({e})")
    out_path = out_path or _default_out(ir, target_lang)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(out_path)
    return out_path


def _extract_rtf(p: Path) -> dict[str, Any]:
    # rtf → markdown, then treated as a text-native document (Path A, ADR-014).
    import pypandoc
    md = pypandoc.convert_file(str(p), "markdown", format="rtf")
    return {"format": "rtf", "source_path": str(p), "text": md}


def _recompose_rtf(ir: dict[str, Any], target_lang: str, out_path: str | None) -> str:
    import pypandoc
    md = ir.get("doc_out", {}).get(target_lang, ir["text"])
    out_path = out_path or _default_out(ir, target_lang)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    pypandoc.convert_text(md, "rtf", format="markdown", outputfile=out_path)
    return out_path


# --- Tier 3: .pptx (experimental, post-MVP — lesson #12, §4.2/§10) ---
# Best-effort. Text runs per shape at paragraph granularity; recomposition of
# complex layouts may degrade. The UI warns before submission.


def _iter_pptx_paragraphs(prs: Any, include_notes: bool = False) -> list[Any]:
    """Every text paragraph in a presentation, stable order: for each slide, each
    shape's text frame, then table cells; and — when ``include_notes`` — the
    slide's speaker-notes paragraphs. extract + recompose walk this identically."""
    paras: list[Any] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                paras.extend(shape.text_frame.paragraphs)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        paras.extend(cell.text_frame.paragraphs)
        if include_notes:
            try:
                if slide.has_notes_slide:
                    paras.extend(slide.notes_slide.notes_text_frame.paragraphs)
            except Exception:
                pass
    return paras


# pptx run-fidelity (ADR-018) — the docx approach ported to DrawingML (a:p/a:r).
# Simpler than docx: no hyperlink/drawing runs inside a text paragraph (images are
# separate shapes). Group a:r runs by a:rPr, mark spans, rebuild keeping each rPr.

def _pptx_run_text(r: Any, qn: Any) -> str:
    t = r.find(qn("a:t"))
    return (t.text or "") if t is not None else ""


def _pptx_rpr_sig(r: Any, qn: Any) -> Any:
    from lxml import etree
    rpr = r.find(qn("a:rPr"))
    return etree.tostring(rpr) if rpr is not None else b""


def _pptx_p_spans(p_el: Any, qn: Any) -> list[dict[str, Any]]:
    spans: list[dict[str, Any]] = []
    for r in p_el.findall(qn("a:r")):
        sig = _pptx_rpr_sig(r, qn)
        if spans and spans[-1]["sig"] == sig:
            spans[-1]["runs"].append(r)
        else:
            spans.append({"sig": sig, "runs": [r]})
    return spans


def _pptx_set_run_text(r: Any, text: str, qn: Any) -> None:
    t = r.find(qn("a:t"))
    if t is None:
        from lxml import etree
        t = etree.SubElement(r, qn("a:t"))
    t.text = text


def _pptx_serialize_p(para: Any, qn: Any) -> tuple[str, bool]:
    spans = _pptx_p_spans(para._p, qn)
    texts = ["".join(_pptx_run_text(r, qn) for r in s["runs"]) for s in spans]
    full = "".join(texts)
    if not full.strip():
        return full, False
    if len(spans) <= 1:
        return full, True
    return "".join(f"⟦{i}⟧{t}⟦/{i}⟧" for i, t in enumerate(texts)), True


def _pptx_rebuild_p(para: Any, translated: str, qn: Any) -> None:
    spans = _pptx_p_spans(para._p, qn)
    if not spans:
        para.text = _strip_markers(translated)  # python-pptx replaces runs with one
        return
    if len(spans) == 1:
        texts: list[str] = [translated]
    else:
        parsed = _parse_markers(translated, len(spans))
        texts = parsed if parsed is not None else [_strip_markers(translated)] + [""] * (len(spans) - 1)
    for span, text in zip(spans, texts):
        runs = span["runs"]
        _pptx_set_run_text(runs[0], text, qn)
        for extra in runs[1:]:
            extra.getparent().remove(extra)


def _extract_pptx(p: Path, include_notes: bool = False) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    prs = Presentation(str(p))
    segments = []
    for para in _iter_pptx_paragraphs(prs, include_notes):
        text, translate = _pptx_serialize_p(para, qn)
        segments.append({"text": text, "translate": translate})
    return {"format": "pptx", "source_path": str(p), "segments": segments, "include_notes": include_notes}


def _recompose_pptx(ir: dict[str, Any], target_lang: str, out_path: str | None) -> str:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    prs = Presentation(ir["source_path"])  # fresh per language
    rtl = _is_rtl(target_lang)
    for seg, para in zip(ir["segments"], _iter_pptx_paragraphs(prs, ir.get("include_notes", False))):
        if not seg["translate"]:
            continue
        translated = seg.get("out", {}).get(target_lang)
        if translated is None:
            continue
        try:
            _pptx_rebuild_p(para, translated, qn)
        except Exception as e:
            logger.warning(f"pptx paragraph rebuild failed ({e}); plain fallback")
            try:
                para.text = _strip_markers(translated)
            except Exception:
                pass
        if rtl:  # DrawingML direction is paragraph-level (a:pPr @rtl), ADR-020
            try:
                para._p.get_or_add_pPr().set("rtl", "1")
            except Exception as e:
                logger.warning(f"pptx paragraph RTL failed ({e})")
    out_path = out_path or _default_out(ir, target_lang)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def _default_out(ir: dict[str, Any], target_lang: str) -> str:
    p = Path(ir["source_path"])
    return str(p.with_name(f"{p.stem}.{target_lang}{p.suffix}"))


# ---------------------------------------------------------------------------
# extract dispatch
# ---------------------------------------------------------------------------


def extract(path: str, options: dict[str, Any] | None = None) -> dict[str, Any]:
    """Build the IR for a document, dispatching on format (Tier 1 + Tier 2).

    ``options`` are per-job document options (§13.2): ``translate_footnotes``
    (docx, default on) and ``translate_speaker_notes`` (pptx, default off)."""
    options = options or {}
    p = Path(path)
    suf = p.suffix.lower()
    if suf in _DOCX_SUFFIXES:
        return _extract_docx(p, options.get("translate_footnotes", True))
    if suf in _RTF_SUFFIXES:
        return _extract_rtf(p)
    if suf in _PPTX_SUFFIXES:
        return _extract_pptx(p, options.get("translate_speaker_notes", False))
    # Text-native (Path A, ADR-014): keep the whole text; the model preserves
    # Markdown structure in one pass (no loss-less segmentation needed).
    fmt = "md" if suf in _MD_SUFFIXES else "txt"
    return {"format": fmt, "source_path": str(p), "text": p.read_text(encoding="utf-8")}


def count_translatable_chars(path: str) -> int:
    """Translatable character count — drives the duration estimate. Text-native
    docs count the whole text; structure-native (docx/pptx) count paragraph text."""
    ir = extract(path)
    if "text" in ir:
        return len(ir["text"])
    return sum(len(s["text"]) for s in ir["segments"] if s["translate"])


# ---------------------------------------------------------------------------
# Per-segment two-pass (draft → glossary review)
# ---------------------------------------------------------------------------


def _engine_config(
    settings: dict[str, Any], connection_id: str | None, model: str | None
) -> dict[str, Any]:
    """The single translation engine (Sprint 9 / ADR-010): connection + model +
    temperature (default 0.1) + enable_thinking. An explicit ``connection_id`` /
    ``model`` overrides the slot but **keeps** the slot temperature + thinking."""
    temp = settings.get("translation_temperature")
    return {
        "connection_id": connection_id or settings.get("translation_connection") or "",
        "model": model or settings.get("translation_model") or "",
        "temperature": 0.1 if temp is None else temp,
        "max_tokens": settings.get("translation_max_tokens"),
        "num_ctx": settings.get("translation_num_ctx"),
        # Hardcoded OFF for translation (not admin-configurable): a reasoning pass
        # drains the output-token budget and — when the visible answer returns empty
        # after <think> stripping — silently blanks a chunk (the "queso-suizo" bug).
        # The adapter maps this intent to each model family's switch.
        "enable_thinking": False,
        "_slot_name": "engine",
    }


def _resolve_chain(
    connection_id: str | None = None, model: str | None = None, frontend_id: str = ""
) -> list[dict[str, Any]]:
    """Single-engine chain, frontend-aware (per-server model/temperature override)."""
    return [_engine_config(get_llm_settings(frontend_id), connection_id, model)]


async def _translate_unit(
    core: str,
    *,
    source_name: str,
    target_name: str,
    ctx_before: str,
    ctx_after: str,
    sliced: list[dict[str, str]],
    system_prompt: str,
    review_prompt: str,
    chain: list[dict[str, Any]],
) -> str:
    """Translate one unit with a completeness guard: retry when the result comes
    back empty/truncated, and fail loudly rather than ship a hole (queso-suizo bug)."""
    last = ""
    for attempt in range(1, _TRANSLATION_MAX_ATTEMPTS + 1):
        last = await _translate_unit_once(
            core, source_name=source_name, target_name=target_name,
            ctx_before=ctx_before, ctx_after=ctx_after, sliced=sliced,
            system_prompt=system_prompt, review_prompt=review_prompt, chain=chain)
        if not _looks_truncated(core, last):
            return last
        logger.warning(
            f"Translation → {target_name} looks incomplete "
            f"({len(last.strip())} vs {len(core.strip())} src chars), "
            f"attempt {attempt}/{_TRANSLATION_MAX_ATTEMPTS} — retrying")
    raise ValueError(
        f"Translation → {target_name} still incomplete after {_TRANSLATION_MAX_ATTEMPTS} "
        f"attempts ({len(last.strip())} vs {len(core.strip())} src chars); failing the job "
        f"rather than emailing a truncated document")


async def _translate_unit_once(
    core: str,
    *,
    source_name: str,
    target_name: str,
    ctx_before: str,
    ctx_after: str,
    sliced: list[dict[str, str]],
    system_prompt: str,
    review_prompt: str,
    chain: list[dict[str, Any]],
) -> str:
    """Translate one unit of text (a whole text-native document/chunk, or one
    structural paragraph). Two independent LLM calls, each with its own
    self-contained prompt (like n8n's separate Translate and Review nodes):
    pass 1 (translate) → pass 2 (glossary review) only when the worklist is
    non-empty (§11.4 — skip the empty pass-2 call). Optional neighbour context."""
    context_parts = []
    if ctx_before.strip():
        context_parts.append(f"[preceding text]\n{ctx_before.strip()}")
    if ctx_after.strip():
        context_parts.append(f"[following text]\n{ctx_after.strip()}")
    context_block = (
        "--- Context (for consistency only — do NOT translate this) ---\n"
        + "\n\n".join(context_parts) + "\n\n"
    ) if context_parts else ""

    pass1_user = (
        f"Source language: {source_name}. Target language: {target_name}.\n\n"
        f"{context_block}--- Text to translate ---\n{core}"
    )
    draft = _strip_think(await llm.chat_with_fallback(
        [{"role": "system", "content": system_prompt},
         {"role": "user", "content": pass1_user}],
        chain,
    ))

    if not sliced:
        return draft  # no glossary worklist → one pass, zero pass-2 call

    review_user = (
        f"Target language: {target_name}.\n\n"
        f"Draft translation:\n{draft}\n\n{format_glossary_block(sliced)}\n\nReturn the corrected text only."
    )
    return _strip_think(await llm.chat_with_fallback(
        [{"role": "system", "content": review_prompt},
         {"role": "user", "content": review_user}],
        chain,
    ))


def _is_heading(seg: dict[str, Any]) -> bool:
    """A translatable segment whose first line opens a markdown heading (``#``)."""
    return bool(seg["translate"]) and seg["text"].lstrip().startswith("#")


def _chunk_text(text: str, target: int) -> list[str]:
    """Split a text-native document into chunks near ``target`` chars, cutting at
    **heading boundaries** (§11.5, ADR-019). A running chunk is closed when the next
    section (heading) would push it past the target, so each chunk starts at a heading
    and stays a semantically-closed unit. An oversized single section (no inner heading)
    falls back to cutting at block boundaries once the chunk reaches the target — never
    mid-paragraph. Loss-less: ``"".join(chunks) == text``."""
    if len(text) <= target:
        return [text]
    chunks: list[str] = []
    cur = ""
    for seg in _segment(text, "md"):
        starts_block = bool(seg["translate"])  # a translatable run begins a block
        would_overflow = bool(cur) and len(cur) + len(seg["text"]) > target
        # Cut before a heading that would overflow (preferred), or before any block
        # once the running chunk already reached the target (oversized-section fallback).
        if starts_block and would_overflow and (_is_heading(seg) or len(cur) >= target):
            chunks.append(cur)
            cur = ""
        cur += seg["text"]
    if cur:
        chunks.append(cur)
    return chunks or [text]


def _split_trailing_newlines(text: str) -> tuple[str, str]:
    """Return (content, trailing_newlines) so recompose keeps block spacing."""
    core = text.rstrip("\n")
    return core, text[len(core):]


# ---------------------------------------------------------------------------
# Translate + recompose
# ---------------------------------------------------------------------------


async def _reading_pass(ir: dict[str, Any], source_name: str, chain: list[dict[str, Any]]) -> str:
    """§13.3 — one reading pass over the document to produce a neutral context brief,
    cached on the IR (computed once, reused per language). Reads the whole text for a
    text-native doc (Path A, when chunked) or the joined translatable segments (Path B)."""
    if "text" in ir:
        text = ir["text"]
    else:
        text = "\n".join(_strip_markers(s["text"]) for s in ir["segments"] if s["translate"])
    text = text[: config.translation_input_budget_chars]
    if not text.strip():
        return ""
    try:
        brief = _strip_think(await llm.chat_with_fallback(
            [{"role": "system", "content": load_context_prompt()},
             {"role": "user", "content": f"Source language: {source_name}.\n\n--- Document ---\n{text}"}],
            chain,
        ))
        logger.info("Reading pass produced a document-context brief")
        return brief.strip()
    except Exception as e:
        logger.warning(f"reading pass failed (translating without a context brief): {e}")
        return ""


async def translate(
    ir: dict[str, Any],
    source_lang: str,
    target_langs: list[str],
    user_glossary: str = "",
    connection_id: str | None = None,
    model: str | None = None,
    frontend_id: str = "",
    options: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Translate ``ir`` into each target language, dispatching by path (ADR-014):
    text-native docs (md/txt/rtf) whole-document per language; structure-native
    docs (docx/pptx) per structural unit as plain text. Mutates and returns ``ir``."""
    options = options or {}
    chain = _resolve_chain(connection_id, model, frontend_id)
    if not chain or not chain[0].get("connection_id"):
        raise ValueError("no translation connection resolved (register one / set the translation slot)")

    # Resolve the effective glossary once per job: base + per-server (replace|append);
    # the per-job user glossary wins per unit (§11.3, precedence user > server > base).
    # When the admin has turned the glossary off ("ignore glossary" toggle, §11.4a),
    # every layer — base, per-server AND the per-job user terms — is skipped.
    from src.api.v1.admin.knowledge import resolve_glossary
    glossary_on = get_llm_settings(frontend_id).get("translation_glossary_enabled", True)
    glossary_terms = resolve_glossary(frontend_id) if glossary_on else []
    if not glossary_on:
        user_glossary = ""
    flavour = load_document_flavour(frontend_id)
    flavour_block = f"\n\n## House style (this instance)\n\n{flavour}" if flavour else ""
    source_name = language_name(source_lang)

    common = dict(
        source_lang=source_lang, target_langs=target_langs, user_glossary=user_glossary,
        glossary_terms=glossary_terms, chain=chain, source_name=source_name,
        review_prompt=load_review_prompt(),  # pass-2 procedure — no flavour, blind to pass 1
    )
    if "text" in ir:  # Path A — text-native: markdown-aware, heading-chunked (§11.5)
        # Once Path A is chunked, each chunk loses whole-document context, so inject the
        # same neutral §13.3 brief. An unchunked document is one call with full context.
        context_block = ""
        if options.get("contextual", True) and len(ir["text"]) > config.translation_chunk_target_chars:
            if "_context" not in ir:
                ir["_context"] = await _reading_pass(ir, source_name, chain)
            if ir.get("_context"):
                context_block = f"\n\n## Document context (for cross-chunk coherence)\n\n{ir['_context']}"
        await _translate_text_native(
            ir, system_prompt=load_document_translation_prompt() + flavour_block + context_block, **common)
    else:  # Path B — structure-native: per unit, plain text
        # Contextual translation (§13.3, default on): a reading pass gives the
        # per-unit calls the whole-document gist. Computed once, cached on the IR.
        context_block = ""
        if options.get("contextual", True):
            if "_context" not in ir:
                ir["_context"] = await _reading_pass(ir, source_name, chain)
            if ir.get("_context"):
                context_block = f"\n\n## Document context (for disambiguating short strings)\n\n{ir['_context']}"
        await _translate_structural(
            ir, system_prompt=load_plain_translation_prompt() + flavour_block + context_block, **common)
    return ir


def _log_coverage(covered: int, lang: str) -> None:
    if covered == 0:
        logger.info(f"No glossary coverage for target '{lang}' — one-pass translation (pass-2 skipped)")


async def _translate_text_native(
    ir: dict[str, Any], *, source_lang, target_langs, user_glossary, glossary_terms,
    chain, source_name, system_prompt, review_prompt,
) -> None:
    """Path A: translate the whole document per language in one call (two-pass);
    split into top-level-block chunks only when the source exceeds the budget."""
    chunks = _chunk_text(ir["text"], config.translation_chunk_target_chars)
    for lang in target_langs:
        target_name = language_name(lang)
        covered, parts = 0, []
        for chunk in chunks:
            core, trailing = _split_trailing_newlines(chunk)
            sliced = slice_glossary(core, source_lang, lang, user_glossary, base_terms=glossary_terms)
            covered += len(sliced)
            translated = await _translate_unit(
                core, source_name=source_name, target_name=target_name,
                ctx_before="", ctx_after="", sliced=sliced,
                system_prompt=system_prompt, review_prompt=review_prompt, chain=chain)
            parts.append(translated.rstrip("\n") + trailing)
        ir.setdefault("doc_out", {})[lang] = "".join(parts)
        _log_coverage(covered, lang)
        logger.info(f"Translated whole document ({len(chunks)} chunk(s)) → {lang}")


async def _translate_structural(
    ir: dict[str, Any], *, source_lang, target_langs, user_glossary, glossary_terms,
    chain, source_name, system_prompt, review_prompt,
) -> None:
    """Path B: translate each structural unit (docx/pptx paragraph) as plain text,
    two-pass; structure is preserved by the library at recompose."""
    segs = ir["segments"]
    trans_idx = [i for i, s in enumerate(segs) if s["translate"]]
    for lang in target_langs:
        target_name = language_name(lang)
        covered = 0
        for pos, i in enumerate(trans_idx):
            core, trailing = _split_trailing_newlines(segs[i]["text"])
            ctx_before = segs[trans_idx[pos - 1]]["text"] if pos > 0 else ""
            ctx_after = segs[trans_idx[pos + 1]]["text"] if pos < len(trans_idx) - 1 else ""
            sliced = slice_glossary(core, source_lang, lang, user_glossary, base_terms=glossary_terms)
            covered += len(sliced)
            translated = await _translate_unit(
                core, source_name=source_name, target_name=target_name,
                ctx_before=ctx_before, ctx_after=ctx_after, sliced=sliced,
                system_prompt=system_prompt, review_prompt=review_prompt, chain=chain)
            segs[i].setdefault("out", {})[lang] = translated.rstrip("\n") + trailing
        _log_coverage(covered, lang)
        logger.info(f"Translated {len(trans_idx)} units → {lang}")


def recompose(ir: dict[str, Any], target_lang: str, out_path: str | None = None) -> str:
    """Rebuild the document for one target language and write it, dispatching on
    format. Non-translatable / untranslated segments fall back to source."""
    fmt = ir["format"]
    if fmt == "docx":
        return _recompose_docx(ir, target_lang, out_path)
    if fmt == "rtf":
        return _recompose_rtf(ir, target_lang, out_path)
    if fmt == "pptx":
        return _recompose_pptx(ir, target_lang, out_path)
    # Path A (txt / md): write the whole translated document.
    result = ir.get("doc_out", {}).get(target_lang, ir["text"])
    out_path = out_path or _default_out(ir, target_lang)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    Path(out_path).write_text(result, encoding="utf-8")
    return out_path


async def translate_document(
    path: str,
    source_lang: str,
    target_langs: list[str],
    user_glossary: str = "",
    out_dir: str | None = None,
    connection_id: str | None = None,
    model: str | None = None,
    frontend_id: str = "",
    options: dict[str, Any] | None = None,
) -> dict[str, str]:
    """End-to-end: extract → translate → recompose. Returns ``{lang: out_path}``."""
    ir = extract(path, options)
    await translate(ir, source_lang, target_langs, user_glossary, connection_id, model, frontend_id)
    src = Path(path)
    outputs: dict[str, str] = {}
    for lang in target_langs:
        out_path = (
            str(Path(out_dir) / f"{src.stem}.{lang}{src.suffix}") if out_dir else None
        )
        outputs[lang] = recompose(ir, lang, out_path)
    return outputs
